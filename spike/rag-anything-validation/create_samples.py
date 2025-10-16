#!/usr/bin/env python3
"""
Generate sample documents for RAG-Anything validation spike.
Creates PDF, DOCX, and PPTX files with Climate Change Science content.
"""

import os
from pathlib import Path

# PDF Generation
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors

    def create_pdf():
        """Create text-based PDF with tables and equations"""
        pdf_path = Path("samples/climate-research-paper.pdf")
        doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()

        # Title
        title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontSize=24, spaceAfter=30)
        story.append(Paragraph("Climate Change: Temperature and CO₂ Correlation Analysis", title_style))
        story.append(Spacer(1, 0.2 * inch))

        # Abstract
        story.append(Paragraph("<b>Abstract</b>", styles['Heading2']))
        abstract_text = """
        This research paper examines the correlation between atmospheric carbon dioxide concentrations
        and global temperature anomalies over the past 140 years. Analysis of historical data demonstrates
        a strong positive correlation (r=0.94) between CO₂ levels and temperature increases, supporting
        the hypothesis of anthropogenic climate forcing.
        """
        story.append(Paragraph(abstract_text, styles['BodyText']))
        story.append(Spacer(1, 0.3 * inch))

        # Introduction with equation
        story.append(Paragraph("<b>1. Mathematical Framework</b>", styles['Heading2']))
        story.append(Paragraph(
            "The radiative forcing equation for CO₂ is approximated as:",
            styles['BodyText']
        ))
        story.append(Spacer(1, 0.1 * inch))

        # LaTeX-style equation representation
        equation_style = ParagraphStyle('Equation', parent=styles['BodyText'],
                                       alignment=1, fontName='Courier')
        story.append(Paragraph("ΔF = 5.35 × ln(C/C₀) W/m²", equation_style))
        story.append(Spacer(1, 0.1 * inch))
        story.append(Paragraph(
            "where C is current CO₂ concentration and C₀ is the reference concentration (280 ppm).",
            styles['BodyText']
        ))
        story.append(Spacer(1, 0.3 * inch))

        # Data table
        story.append(Paragraph("<b>2. Observational Data</b>", styles['Heading2']))
        story.append(Paragraph("Table 1: Key Climate Indicators by Decade", styles['BodyText']))
        story.append(Spacer(1, 0.1 * inch))

        data = [
            ['Decade', 'Temp Anomaly (°C)', 'CO₂ (ppm)', 'Sea Level Rise (mm)'],
            ['1880s', '-0.16', '280', '0'],
            ['1980s', '0.26', '339', '80'],
            ['2020s', '1.27', '420', '230']
        ]

        table = Table(data, colWidths=[1.5*inch, 1.5*inch, 1.2*inch, 1.5*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(table)
        story.append(Spacer(1, 0.3 * inch))

        # Results
        story.append(Paragraph("<b>3. Results and Discussion</b>", styles['Heading2']))
        results_text = """
        Linear regression analysis reveals a strong correlation between CO₂ concentrations and
        temperature anomalies (r² = 0.88, p < 0.001). The temperature sensitivity to CO₂ doubling
        (climate sensitivity) is estimated at 3.0°C ± 0.5°C, consistent with IPCC projections.
        """
        story.append(Paragraph(results_text, styles['BodyText']))
        story.append(Spacer(1, 0.2 * inch))

        # Conclusion
        story.append(Paragraph("<b>4. Conclusions</b>", styles['Heading2']))
        conclusion_text = """
        This analysis provides quantitative evidence for the link between anthropogenic CO₂ emissions
        and observed global warming. The findings underscore the urgency of emissions reduction to
        limit temperature increases to 1.5°C above pre-industrial levels.
        """
        story.append(Paragraph(conclusion_text, styles['BodyText']))

        doc.build(story)
        print(f"✓ Created: {pdf_path}")

except ImportError:
    print("⚠ reportlab not installed, skipping PDF generation")


# DOCX Generation
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    def create_docx():
        """Create Word document with tables and images placeholder"""
        doc = Document()

        # Title
        title = doc.add_heading('Climate Change Mitigation Strategies', level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Introduction
        doc.add_heading('1. Introduction', level=1)
        doc.add_paragraph(
            'Climate change mitigation encompasses actions to limit the magnitude and rate of '
            'long-term climate change. This document outlines evidence-based strategies for '
            'reducing greenhouse gas emissions and enhancing carbon sinks.'
        )

        # Mitigation Strategies with table
        doc.add_heading('2. Mitigation Strategy Matrix', level=1)
        doc.add_paragraph('Table 1: Priority Mitigation Actions by Sector')

        table = doc.add_table(rows=5, cols=4)
        table.style = 'Light Grid Accent 1'

        # Header row
        header_cells = table.rows[0].cells
        header_cells[0].text = 'Sector'
        header_cells[1].text = 'Strategy'
        header_cells[2].text = 'Potential (GtCO₂/yr)'
        header_cells[3].text = 'Implementation Cost'

        # Data rows
        data_rows = [
            ['Energy', 'Renewable transition', '12.5', 'Medium'],
            ['Transport', 'Electrification', '4.2', 'High'],
            ['Buildings', 'Efficiency improvements', '3.8', 'Low'],
            ['Agriculture', 'Sustainable practices', '5.1', 'Low']
        ]

        for i, row_data in enumerate(data_rows, start=1):
            cells = table.rows[i].cells
            for j, value in enumerate(row_data):
                cells[j].text = value

        doc.add_paragraph()  # Spacer

        # Renewable Energy Section
        doc.add_heading('3. Renewable Energy Transition', level=1)
        doc.add_paragraph(
            'Transitioning from fossil fuels to renewable energy sources represents the largest '
            'mitigation opportunity, with potential to reduce emissions by 12+ gigatons annually.'
        )

        doc.add_heading('Key Technologies:', level=2)
        doc.add_paragraph('Solar photovoltaic systems', style='List Bullet')
        doc.add_paragraph('Wind energy (onshore and offshore)', style='List Bullet')
        doc.add_paragraph('Hydroelectric power', style='List Bullet')
        doc.add_paragraph('Geothermal energy', style='List Bullet')
        doc.add_paragraph('Green hydrogen production', style='List Bullet')

        # Image placeholder
        doc.add_heading('4. Global Renewable Capacity Growth', level=1)
        doc.add_paragraph('[IMAGE: Chart showing renewable energy capacity 2000-2023]')
        doc.add_paragraph('Note: Renewable energy capacity has grown exponentially, '
                         'with solar and wind leading the expansion.')

        # Carbon Sequestration
        doc.add_heading('5. Carbon Sequestration Methods', level=1)
        doc.add_paragraph(
            'Beyond emissions reduction, removing CO₂ from the atmosphere through natural and '
            'technological means is essential for achieving net-zero targets.'
        )

        doc.add_paragraph('Afforestation and reforestation programs', style='List Number')
        doc.add_paragraph('Soil carbon enhancement through regenerative agriculture', style='List Number')
        doc.add_paragraph('Direct air capture (DAC) technologies', style='List Number')
        doc.add_paragraph('Ocean-based sequestration (blue carbon)', style='List Number')

        # Conclusion
        doc.add_heading('6. Conclusion', level=1)
        doc.add_paragraph(
            'Achieving the 1.5°C target requires rapid deployment of mitigation strategies across '
            'all economic sectors. The combination of renewable energy expansion, energy efficiency, '
            'and carbon sequestration can limit warming while supporting sustainable development goals.'
        )

        docx_path = Path("samples/climate-mitigation-strategies.docx")
        doc.save(str(docx_path))
        print(f"✓ Created: {docx_path}")

except ImportError:
    print("⚠ python-docx not installed, skipping DOCX generation")


# PPTX Generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    def create_pptx():
        """Create PowerPoint presentation with slide layouts and charts placeholder"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Climate Change Science"
        subtitle.text = "Understanding Global Warming Trends and Impacts\nOctober 2023"

        # Slide 2: Overview
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Presentation Overview"
        tf = body_shape.text_frame
        tf.text = "Key topics covered in this presentation:"

        for point in [
            "Temperature trends and observations",
            "Carbon dioxide concentration increases",
            "Impacts on ecosystems and human systems",
            "Future projections and scenarios",
            "Mitigation and adaptation strategies"
        ]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

        # Slide 3: Temperature Data (with table)
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        shapes = slide.shapes

        title_shape = shapes.title
        title_shape.text = "Global Temperature Anomalies"

        # Add table
        rows, cols = 5, 3
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(7.0)
        height = Inches(3.0)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2.5)
        table.columns[2].width = Inches(2.0)

        # Header row
        table.cell(0, 0).text = "Period"
        table.cell(0, 1).text = "Temperature Anomaly (°C)"
        table.cell(0, 2).text = "Trend"

        # Data rows
        data = [
            ("1880-1920", "+0.05", "Stable"),
            ("1920-1960", "+0.10", "Slight increase"),
            ("1960-2000", "+0.45", "Accelerating"),
            ("2000-2023", "+1.20", "Rapid warming")
        ]

        for i, (period, temp, trend) in enumerate(data, start=1):
            table.cell(i, 0).text = period
            table.cell(i, 1).text = temp
            table.cell(i, 2).text = trend

        # Slide 4: CO2 Trends (with chart placeholder)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes

        title_shape = shapes.title
        title_shape.text = "Atmospheric CO₂ Concentration"

        # Chart placeholder text
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "[CHART: CO₂ Concentration 1880-2023]"

        p = tf.add_paragraph()
        p.text = "Key observations:"
        p = tf.add_paragraph()
        p.text = "• Pre-industrial: 280 ppm"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• Current (2023): 420 ppm"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "• Rate of increase: 2-3 ppm/year"
        p.level = 1

        # Slide 5: Impacts
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Observed Climate Impacts"
        tf = body_shape.text_frame
        tf.text = "Arctic sea ice decline (13% per decade)"

        impacts = [
            "Sea level rise (3.7 mm/year)",
            "Increased extreme weather frequency",
            "Ocean acidification and warming",
            "Ecosystem disruption and species migration",
            "Agricultural productivity threats"
        ]

        for impact in impacts:
            p = tf.add_paragraph()
            p.text = impact
            p.level = 0

        # Slide 6: Future Scenarios
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Future Temperature Projections"
        tf = body_shape.text_frame
        tf.text = "Warming by 2100 depends on emissions pathway:"

        scenarios = [
            ("RCP 2.6 (Strong mitigation)", "1.5-2.0°C warming"),
            ("RCP 4.5 (Moderate action)", "2.0-3.0°C warming"),
            ("RCP 8.5 (Business as usual)", "3.5-5.0°C warming")
        ]

        for scenario, result in scenarios:
            p = tf.add_paragraph()
            p.text = f"{scenario}: {result}"
            p.level = 1

        # Slide 7: Conclusion
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Key Takeaways"
        tf = body_shape.text_frame
        tf.text = "Climate change is unequivocal and primarily anthropogenic"

        takeaways = [
            "Temperature has increased 1.1°C since pre-industrial era",
            "CO₂ levels are highest in 3 million years",
            "Impacts are already visible across multiple systems",
            "Limiting warming to 1.5°C requires immediate action",
            "Both mitigation and adaptation are essential"
        ]

        for takeaway in takeaways:
            p = tf.add_paragraph()
            p.text = takeaway
            p.level = 0

        pptx_path = Path("samples/climate-science-presentation.pptx")
        prs.save(str(pptx_path))
        print(f"✓ Created: {pptx_path}")

except ImportError:
    print("⚠ python-pptx not installed, skipping PPTX generation")


def create_scanned_pdf_note():
    """Create note about scanned PDF requirement"""
    note_path = Path("samples/SCANNED_PDF_NOTE.txt")
    note_path.write_text(
        "SCANNED PDF REQUIREMENT:\n\n"
        "A scanned/image-based PDF is needed for OCR testing.\n"
        "Since we cannot easily generate this programmatically,\n"
        "options:\n"
        "1. Print one of the generated PDFs and scan it\n"
        "2. Convert a page to image and embed in PDF\n"
        "3. Download a sample scanned document\n"
        "4. Use tesseract to generate one from an image\n\n"
        "For spike purposes, we can test OCR capabilities on\n"
        "a converted image-based PDF or skip this specific test\n"
        "and note it in the spike report.\n"
    )
    print(f"✓ Created: {note_path}")


if __name__ == "__main__":
    print("Generating sample documents for RAG-Anything validation...\n")

    # Ensure samples directory exists
    Path("samples").mkdir(exist_ok=True)

    create_pdf()
    create_docx()
    create_pptx()
    create_scanned_pdf_note()

    print("\n✓ Sample document generation complete!")
    print("\nNote: TXT, MD, and CSV files should already exist in samples/")
