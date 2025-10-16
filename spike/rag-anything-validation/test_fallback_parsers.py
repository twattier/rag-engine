#!/usr/bin/env python3
"""
Test fallback parsers (python-docx, pypdf, pandas) as alternative to RAG-Anything.
This provides immediate spike validation while RAG-Anything installation completes.

Story 0.1: RAG-Anything Technical Validation Spike
"""

from __future__ import annotations

import json
import time
from pathlib import Path
from typing import Any
import sys

# Track results
results = {
    "formats_tested": [],
    "successful_formats": [],
    "failed_formats": [],
    "parsing_times": {},
    "errors": {},
    "output_samples": {}
}


def test_txt_format(file_path: Path) -> dict[str, Any]:
    """Test plain text parsing"""
    result = {
        "format": "Plain Text",
        "file": str(file_path),
        "success": False,
        "parse_time_seconds": 0.0,
        "error": None,
        "text_sample": None
    }

    try:
        start_time = time.perf_counter()
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        end_time = time.perf_counter()

        result["parse_time_seconds"] = end_time - start_time
        result["text_sample"] = content[:200] + "..." if len(content) > 200 else content
        result["success"] = True

    except Exception as e:
        result["error"] = f"{type(e).__name__}: {str(e)}"

    return result


def test_md_format(file_path: Path) -> dict[str, Any]:
    """Test markdown parsing"""
    result = {
        "format": "Markdown",
        "file": str(file_path),
        "success": False,
        "parse_time_seconds": 0.0,
        "error": None,
        "text_sample": None
    }

    try:
        start_time = time.perf_counter()
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        end_time = time.perf_counter()

        result["parse_time_seconds"] = end_time - start_time
        result["text_sample"] = content[:200] + "..." if len(content) > 200 else content
        result["success"] = True

    except Exception as e:
        result["error"] = f"{type(e).__name__}: {str(e)}"

    return result


def test_csv_format(file_path: Path) -> dict[str, Any]:
    """Test CSV parsing with pandas"""
    result = {
        "format": "CSV",
        "file": str(file_path),
        "success": False,
        "parse_time_seconds": 0.0,
        "error": None,
        "text_sample": None
    }

    try:
        import pandas as pd

        start_time = time.perf_counter()
        df = pd.read_csv(file_path)
        end_time = time.perf_counter()

        result["parse_time_seconds"] = end_time - start_time
        result["text_sample"] = df.head(5).to_string()
        result["success"] = True

    except ImportError:
        result["error"] = "pandas not installed"
    except Exception as e:
        result["error"] = f"{type(e).__name__}: {str(e)}"

    return result


def test_docx_format(file_path: Path) -> dict[str, Any]:
    """Test DOCX parsing with python-docx"""
    result = {
        "format": "Word DOCX",
        "file": str(file_path),
        "success": False,
        "parse_time_seconds": 0.0,
        "error": None,
        "text_sample": None
    }

    try:
        from docx import Document

        start_time = time.perf_counter()
        doc = Document(file_path)

        # Extract text from paragraphs
        paragraphs = [para.text for para in doc.paragraphs]
        full_text = "\n".join(paragraphs)

        # Extract tables
        tables_found = len(doc.tables)

        end_time = time.perf_counter()

        result["parse_time_seconds"] = end_time - start_time
        result["text_sample"] = full_text[:200] + "..." if len(full_text) > 200 else full_text
        result["tables_found"] = tables_found
        result["success"] = True

    except ImportError:
        result["error"] = "python-docx not installed"
    except Exception as e:
        result["error"] = f"{type(e).__name__}: {str(e)}"

    return result


def test_pptx_format(file_path: Path) -> dict[str, Any]:
    """Test PPTX parsing with python-pptx"""
    result = {
        "format": "PowerPoint PPTX",
        "file": str(file_path),
        "success": False,
        "parse_time_seconds": 0.0,
        "error": None,
        "text_sample": None
    }

    try:
        from pptx import Presentation

        start_time = time.perf_counter()
        prs = Presentation(file_path)

        # Extract text from all slides
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)

        full_text = "\n".join(text_runs)

        end_time = time.perf_counter()

        result["parse_time_seconds"] = end_time - start_time
        result["text_sample"] = full_text[:200] + "..." if len(full_text) > 200 else full_text
        result["slides_found"] = len(prs.slides)
        result["success"] = True

    except ImportError:
        result["error"] = "python-pptx not installed"
    except Exception as e:
        result["error"] = f"{type(e).__name__}: {str(e)}"

    return result


def test_pdf_format(file_path: Path) -> dict[str, Any]:
    """Test PDF parsing with pypdf"""
    result = {
        "format": "PDF",
        "file": str(file_path),
        "success": False,
        "parse_time_seconds": 0.0,
        "error": None,
        "text_sample": None
    }

    try:
        from pypdf import PdfReader

        start_time = time.perf_counter()
        reader = PdfReader(file_path)

        # Extract text from all pages
        text_pages = []
        for page in reader.pages:
            text_pages.append(page.extract_text())

        full_text = "\n".join(text_pages)

        end_time = time.perf_counter()

        result["parse_time_seconds"] = end_time - start_time
        result["text_sample"] = full_text[:200] + "..." if len(full_text) > 200 else full_text
        result["pages_found"] = len(reader.pages)
        result["success"] = True

    except ImportError:
        result["error"] = "pypdf not installed - attempting install"
        # Try to install pypdf
        import subprocess
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "pypdf"])
            # Retry after install
            return test_pdf_format(file_path)
        except Exception as install_err:
            result["error"] = f"Failed to install pypdf: {install_err}"
    except Exception as e:
        result["error"] = f"{type(e).__name__}: {str(e)}"

    return result


def main():
    """Main test execution"""
    print("="*70)
    print("FALLBACK PARSER VALIDATION SPIKE")
    print("Story 0.1: Testing fallback parsers as alternative to RAG-Anything")
    print("="*70 + "\n")

    # Define test files
    samples_dir = Path("samples")
    test_cases = [
        ("Plain Text", samples_dir / "climate-abstract.txt", test_txt_format),
        ("Markdown", samples_dir / "climate-report.md", test_md_format),
        ("CSV", samples_dir / "climate-data.csv", test_csv_format),
        ("Word DOCX", samples_dir / "climate-mitigation-strategies.docx", test_docx_format),
        ("PowerPoint PPTX", samples_dir / "climate-science-presentation.pptx", test_pptx_format),
        ("PDF", samples_dir / "climate-research-paper.pdf", test_pdf_format),
    ]

    # Test each format
    for format_name, file_path, test_func in test_cases:
        print(f"\nTesting: {format_name}")
        print(f"File: {file_path}")
        print("-" * 70)

        if not file_path.exists():
            print(f"⚠ File not found - skipping")
            results["formats_tested"].append(format_name)
            results["failed_formats"].append(format_name)
            results["errors"][format_name] = "File not found"
            continue

        result = test_func(file_path)
        results["formats_tested"].append(format_name)

        if result["success"]:
            results["successful_formats"].append(format_name)
            results["parsing_times"][format_name] = result["parse_time_seconds"]
            results["output_samples"][format_name] = result.get("text_sample", "N/A")

            print(f"✓ Success! Parsed in {result['parse_time_seconds']:.4f}s")
            if "pages_found" in result:
                print(f"  Pages: {result['pages_found']}")
            if "tables_found" in result:
                print(f"  Tables: {result['tables_found']}")
            if "slides_found" in result:
                print(f"  Slides: {result['slides_found']}")
        else:
            results["failed_formats"].append(format_name)
            results["errors"][format_name] = result["error"]
            print(f"✗ Failed: {result['error']}")

    # Save results
    output_dir = Path("outputs")
    output_dir.mkdir(exist_ok=True)

    results_file = output_dir / "fallback_parser_results.json"
    with open(results_file, 'w') as f:
        json.dump(results, f, indent=2)

    print("\n" + "="*70)
    print("SUMMARY")
    print("="*70)
    print(f"Total formats tested: {len(results['formats_tested'])}")
    print(f"Successful: {len(results['successful_formats'])}")
    print(f"Failed: {len(results['failed_formats'])}")
    print(f"\nResults saved to: {results_file}")

    # Print success/failure details
    if results["successful_formats"]:
        print("\n✓ Successful formats:")
        for fmt in results["successful_formats"]:
            time_str = f"{results['parsing_times'][fmt]:.4f}s"
            print(f"  - {fmt}: {time_str}")

    if results["failed_formats"]:
        print("\n✗ Failed formats:")
        for fmt in results["failed_formats"]:
            print(f"  - {fmt}: {results['errors'].get(fmt, 'Unknown error')}")

    # Exit with appropriate code
    success_rate = len(results['successful_formats']) / len(results['formats_tested']) if results['formats_tested'] else 0
    if success_rate >= 0.83:  # 5/6 formats
        print("\n✓ SPIKE VALIDATION: PASS (5+ formats working)")
        sys.exit(0)
    else:
        print(f"\n⚠ SPIKE VALIDATION: PARTIAL ({len(results['successful_formats'])}/{len(results['formats_tested'])} formats working)")
        sys.exit(1)


if __name__ == "__main__":
    main()
