"""Microsoft PowerPoint parser using python-pptx."""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation

from app.parsers.base_parser import BaseParser
from app.models import ContentItem


class PPTXParser(BaseParser):
    """Parser for Microsoft PowerPoint (.pptx) files."""

    async def parse(self, file_path: Path) -> list[ContentItem]:
        """Parse PPTX file and return content."""
        prs = Presentation(file_path)
        content_list: list[ContentItem] = []

        for slide_idx, slide in enumerate(prs.slides):
            # Extract text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                content_list.append(
                    ContentItem(
                        type="text",
                        text="\n".join(slide_text),
                        page_idx=slide_idx,
                        structure=f"pptx_slide_{slide_idx}",
                    )
                )

            # Extract tables from slide
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows.append(row_data)

                    if rows:
                        content_list.append(
                            ContentItem(
                                type="table",
                                rows=rows,
                                page_idx=slide_idx,
                                structure=f"pptx_slide_{slide_idx}_table",
                            )
                        )

        return content_list
